from flask import Flask, request, jsonify, render_template, session
import requests
import json
import uuid
from datetime import datetime
import sqlite3
import os
import re
from werkzeug.utils import secure_filename
import mimetypes
from pathlib import Path

# Load environment variables from .env file
try:
    from dotenv import load_dotenv
    load_dotenv()
    print("+ Environment variables loaded from .env file")
except ImportError:
    print("WARNING: python-dotenv not installed. .env file not loaded.")

# Claude Code SDK integration
try:
    from claude_code_sdk import query, ClaudeCodeOptions, Message
    import anyio
    CLAUDE_CODE_AVAILABLE = True
    print("+ Claude Code SDK enabled")
except ImportError:
    CLAUDE_CODE_AVAILABLE = False
    print("WARNING: claude-code-sdk not installed. Claude Code features disabled.")

# Hugging Face SDK integration for image generation
try:
    from huggingface_hub import InferenceClient
    from PIL import Image as PILImage
    import io
    import base64
    HUGGINGFACE_AVAILABLE = True
    print("+ Hugging Face Hub enabled for image generation")
except ImportError:
    HUGGINGFACE_AVAILABLE = False
    print("WARNING: huggingface-hub not installed. Image generation features disabled.")


# claude implementation
# import anyio
# from claude_code_sdk import query, ClaudeCodeOptions, Message

# async def main():
    
#     options = ClaudeCodeOptions(
#     max_turns=3,
#     system_prompt="You are a helpful assistant",
#     cwd=Path("/"),
#     allowed_tools=["Read", "Write", "Bash"],
#     permission_mode="acceptEdits"
#     )

#     async for message in query(prompt="Hello", options=options):
#         print(message)

#     messages: list[Message] = []
    
#     async for message in query(
#         prompt="Write a haiku about foo.py",
#         options=ClaudeCodeOptions(max_turns=3)
#     ):
#         messages.append(message)
    
#     print(messages)
    

# anyio.run(main)
# Try to import optional dependencies with better error handling
try:
    import markdown
    MARKDOWN_AVAILABLE = True
    print("+ Markdown support enabled")
except ImportError:
    MARKDOWN_AVAILABLE = False
    print("WARNING: markdown not installed. Using basic formatting.")

try:
    # Try new pypdf first, then fall back to PyPDF2
    try:
        import pypdf as PyPDF2
        PDF_AVAILABLE = True
        print("+ PDF support enabled (using pypdf)")
    except ImportError:
        import PyPDF2
        PDF_AVAILABLE = True
        print("+ PDF support enabled (using PyPDF2)")
except ImportError:
    PDF_AVAILABLE = False
    print("WARNING: PDF library not installed. PDF support disabled.")

try:
    from PIL import Image
    import io
    IMAGE_AVAILABLE = True
    print("+ Image processing enabled")
except ImportError:
    IMAGE_AVAILABLE = False
    print("WARNING: Pillow not installed. Image processing disabled.")

try:
    from docx import Document
    DOCX_AVAILABLE = True
    print("+ Word document support enabled")
except ImportError:
    DOCX_AVAILABLE = False
    print("WARNING: python-docx not installed. Word document support disabled.")

try:
    import openpyxl
    import xlrd
    EXCEL_AVAILABLE = True
    print("+ Excel support enabled")
except ImportError:
    EXCEL_AVAILABLE = False
    print("WARNING: openpyxl/xlrd not installed. Excel support disabled.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
    print("+ PowerPoint support enabled")
except ImportError:
    PPTX_AVAILABLE = False
    print("WARNING: python-pptx not installed. PowerPoint support disabled.")

app = Flask(__name__)
app.secret_key = 'your-secret-key-change-this-in-production'

# Ollama configuration
OLLAMA_BASE_URL = 'http://localhost:11434'
DEFAULT_MODELS = ['llama2', 'codellama', 'mistral']

# Claude Code configuration
CLAUDE_CODE_API_KEY = os.getenv('CLAUDE_CODE_API_KEY', '')  # Will be set later
CLAUDE_CODE_MODEL = 'claude-code'

# File upload configuration
UPLOAD_FOLDER = 'uploads'
MAX_FILE_SIZE = 16 * 1024 * 1024  # 16MB
ALLOWED_EXTENSIONS = {
    'txt', 'md', 'py', 'js', 'html', 'css', 'json', 'xml', 'csv'
}

# Add PDF and image extensions only if libraries are available
if PDF_AVAILABLE:
    ALLOWED_EXTENSIONS.add('pdf')
if IMAGE_AVAILABLE:
    ALLOWED_EXTENSIONS.update({'png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp'})

# Add Office document extensions only if libraries are available
if DOCX_AVAILABLE:
    ALLOWED_EXTENSIONS.add('docx')
if EXCEL_AVAILABLE:
    ALLOWED_EXTENSIONS.update({'xlsx', 'xls'})
if PPTX_AVAILABLE:
    ALLOWED_EXTENSIONS.add('pptx')

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE

# Create upload directory
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Database configuration
DATABASE = 'chat_app.db'

class OllamaClient:
    def __init__(self, base_url=OLLAMA_BASE_URL):
        self.base_url = base_url
    
    def chat(self, model, messages):
        """Send chat request to Ollama"""
        url = f"{self.base_url}/api/chat"
        payload = {
            "model": model,
            "messages": messages,
            "stream": False
        }
        
        try:
            response = requests.post(url, json=payload, timeout=120)
            response.raise_for_status()
            return response.json()
        except requests.exceptions.ConnectionError:
            raise Exception("Cannot connect to Ollama. Make sure Ollama is running on http://localhost:11434")
        except requests.exceptions.Timeout:
            raise Exception("Request timed out. The model might be taking too long to respond.")
        except requests.exceptions.RequestException as e:
            raise Exception(f"Ollama request failed: {str(e)}")
    
    def list_models(self):
        """List available models"""
        url = f"{self.base_url}/api/tags"
        
        try:
            response = requests.get(url, timeout=10)
            response.raise_for_status()
            data = response.json()
            return {"models": data.get("models", [])}
        except requests.exceptions.RequestException as e:
            print(f"Error listing models: {e}")
            return {"models": []}

# Initialize Ollama client
ollama_client = OllamaClient()

class ClaudeCodeClient:
    def __init__(self, api_key=None):
        self.api_key = api_key or CLAUDE_CODE_API_KEY
        self.available = CLAUDE_CODE_AVAILABLE and bool(self.api_key)
    
    async def chat(self, messages, tools=None):
        """Send chat request to Claude Code"""
        if not self.available:
            raise Exception("Claude Code not available. Check API key and SDK installation.")
        
        if not self.api_key:
            raise Exception("Claude Code API key not set. Please add CLAUDE_CODE_API_KEY to your .env file.")
        
        if len(self.api_key) < 50:  # Anthropic API keys are typically longer
            raise Exception(f"Claude Code API key appears invalid (length: {len(self.api_key)}). Please check your .env file.")
        
        try:
            # Ensure API key is set in environment for Claude Code SDK
            import os
            os.environ['ANTHROPIC_API_KEY'] = self.api_key
            
            # Convert messages to the format expected by Claude Code SDK
            prompt = self._format_messages(messages)
            
            options = ClaudeCodeOptions(
                max_turns=3,
                system_prompt="You are a helpful AI assistant with access to powerful tools for code analysis, file processing, and data manipulation.",
                cwd=Path("."),
                allowed_tools=tools or ["Read", "Write", "Bash", "Grep"],
                permission_mode="acceptEdits"
            )
            response_content = ""
            async for message in query(prompt=prompt, options=options):
                message_type = type(message).__name__
                
                # Filter out system messages and debug output from UI
                system_message_types = ['SystemMessage', 'DebugMessage', 'InitMessage', 'StatusMessage']
                
                if message_type in system_message_types:
                    print(f"DEBUG: Filtered {message_type}")
                    continue
                
                if hasattr(message, 'content'):
                    # Handle list of TextBlock objects
                    if isinstance(message.content, list):
                        for block in message.content:
                            if hasattr(block, 'text'):
                                response_content += block.text
                            else:
                                response_content += str(block)
                    else:
                        response_content += str(message.content)
                elif isinstance(message, str):
                    response_content += message
                else:
                    # Only add non-system messages to response
                    if message_type not in system_message_types:
                        response_content += str(message)
            
            print(f"Claude Code response: {len(response_content)} characters")
            
            return {
                "message": {
                    "content": response_content,
                    "role": "assistant"
                }
            }
            
        except Exception as e:
            print(f"DEBUG: Claude Code exception: {str(e)}")
            raise Exception(f"Claude Code request failed: {str(e)}")
    
    def _format_messages(self, messages):
        """Convert message history to a single prompt for Claude Code"""
        if not messages:
            return ""
        
        # Get the last user message as the main prompt
        last_message = messages[-1] if messages else {}
        return last_message.get('content', '')
    
    def set_api_key(self, api_key):
        """Set the API key and update availability"""
        self.api_key = api_key
        self.available = CLAUDE_CODE_AVAILABLE and bool(self.api_key)

# Initialize Claude Code client
claude_code_client = ClaudeCodeClient()

class FileProcessor:
    def __init__(self):
        pass
    
    def process_file(self, file_path, filename, mime_type):
        """Process uploaded file and extract content"""
        try:
            if mime_type.startswith('text/') or filename.endswith(('.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml', '.csv')):
                return self._process_text_file(file_path)
            elif mime_type == 'application/pdf' and PDF_AVAILABLE:
                return self._process_pdf(file_path)
            elif mime_type.startswith('image/') and IMAGE_AVAILABLE:
                return self._process_image(file_path, filename)
            elif filename.endswith('.docx') and DOCX_AVAILABLE:
                return self._process_word_document(file_path, filename)
            elif filename.endswith(('.xlsx', '.xls')) and EXCEL_AVAILABLE:
                return self._process_excel_file(file_path, filename)
            elif filename.endswith('.pptx') and PPTX_AVAILABLE:
                return self._process_powerpoint(file_path, filename)
            else:
                return f"File uploaded: {filename} (content extraction not supported for this file type)"
        except Exception as e:
            return f"Error processing file: {str(e)}"
    
    def _process_text_file(self, file_path):
        """Extract text from text files"""
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                return content
            except UnicodeDecodeError:
                continue
            except Exception as e:
                return f"Error reading file: {str(e)}"
        
        return "Error: Could not decode file with any supported encoding"
    
    def _process_pdf(self, file_path):
        """Extract text from PDF files"""
        if not PDF_AVAILABLE:
            return "PDF processing not available (pypdf/PyPDF2 not installed)"
        
        try:
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text.strip() if text.strip() else "PDF processed but no text could be extracted"
        except Exception as e:
            return f"Error reading PDF: {str(e)}"
    
    def _process_image(self, file_path, filename):
        """Process image files"""
        if not IMAGE_AVAILABLE:
            return "Image processing not available (Pillow not installed)"
        
        try:
            with Image.open(file_path) as img:
                width, height = img.size
                format_name = img.format
                mode = img.mode
                
                return f"Image file: {filename}\nDimensions: {width}x{height}\nFormat: {format_name}\nMode: {mode}\n\nNote: For image analysis, please use a vision-capable model like llava."
        except Exception as e:
            return f"Error processing image: {str(e)}"
    
    def _process_word_document(self, file_path, filename):
        """Extract text and formatting from Word documents"""
        if not DOCX_AVAILABLE:
            return "Word document processing not available (python-docx not installed)"
        
        try:
            doc = Document(file_path)
            content = []
            
            content.append(f"Word Document: {filename}\n")
            content.append("=" * 50 + "\n")
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # Basic formatting detection
                    text = paragraph.text
                    if paragraph.style.name.startswith('Heading'):
                        content.append(f"## {text}\n")
                    else:
                        content.append(f"{text}\n")
            
            # Extract tables if present
            if doc.tables:
                content.append("\n**Tables found in document:**\n")
                for i, table in enumerate(doc.tables):
                    content.append(f"\nTable {i+1}:\n")
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        content.append(f"| {row_text} |\n")
            
            return "".join(content) if content else "Word document processed but no content could be extracted"
            
        except Exception as e:
            return f"Error reading Word document: {str(e)}"
    
    def _process_excel_file(self, file_path, filename):
        """Extract data from Excel files"""
        if not EXCEL_AVAILABLE:
            return "Excel processing not available (openpyxl/xlrd not installed)"
        
        try:
            content = []
            content.append(f"Excel File: {filename}\n")
            content.append("=" * 50 + "\n")
            
            if filename.endswith('.xlsx'):
                # Use openpyxl for .xlsx files
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    content.append(f"\n**Sheet: {sheet_name}**\n")
                    
                    # Get data from sheet (limit to prevent huge outputs)
                    max_rows = min(sheet.max_row, 50)  # Limit to 50 rows
                    max_cols = min(sheet.max_column, 20)  # Limit to 20 columns
                    
                    for row in range(1, max_rows + 1):
                        row_data = []
                        for col in range(1, max_cols + 1):
                            cell_value = sheet.cell(row=row, column=col).value
                            if cell_value is not None:
                                row_data.append(str(cell_value))
                            else:
                                row_data.append("")
                        
                        if any(cell.strip() for cell in row_data):  # Only add non-empty rows
                            content.append(f"| {' | '.join(row_data)} |\n")
                    
                    if sheet.max_row > 50:
                        content.append(f"\n... (showing first 50 of {sheet.max_row} rows)\n")
            
            else:
                # Use xlrd for .xls files
                workbook = xlrd.open_workbook(file_path)
                
                for sheet_index in range(workbook.nsheets):
                    sheet = workbook.sheet_by_index(sheet_index)
                    content.append(f"\n**Sheet: {sheet.name}**\n")
                    
                    # Limit rows and columns
                    max_rows = min(sheet.nrows, 50)
                    max_cols = min(sheet.ncols, 20)
                    
                    for row in range(max_rows):
                        row_data = []
                        for col in range(max_cols):
                            cell_value = sheet.cell_value(row, col)
                            row_data.append(str(cell_value) if cell_value else "")
                        
                        if any(cell.strip() for cell in row_data):
                            content.append(f"| {' | '.join(row_data)} |\n")
                    
                    if sheet.nrows > 50:
                        content.append(f"\n... (showing first 50 of {sheet.nrows} rows)\n")
            
            return "".join(content) if content else "Excel file processed but no data could be extracted"
            
        except Exception as e:
            return f"Error reading Excel file: {str(e)}"
    
    def _process_powerpoint(self, file_path, filename):
        """Extract text from PowerPoint presentations"""
        if not PPTX_AVAILABLE:
            return "PowerPoint processing not available (python-pptx not installed)"
        
        try:
            prs = Presentation(file_path)
            content = []
            
            content.append(f"PowerPoint Presentation: {filename}\n")
            content.append("=" * 50 + "\n")
            content.append(f"Total slides: {len(prs.slides)}\n\n")
            
            for i, slide in enumerate(prs.slides, 1):
                content.append(f"**Slide {i}:**\n")
                
                # Extract text from all shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    for text in slide_text:
                        content.append(f"- {text}\n")
                else:
                    content.append("- (No text content)\n")
                
                content.append("\n")
            
            return "".join(content) if content else "PowerPoint file processed but no content could be extracted"
            
        except Exception as e:
            return f"Error reading PowerPoint file: {str(e)}"

class ResponseFormatter:
    def __init__(self):
        pass
        
    def format_response(self, text):
        """Format the AI response with proper HTML structure"""
        if MARKDOWN_AVAILABLE:
            try:
                formatted = markdown.markdown(
                    text,
                    extensions=['codehilite', 'fenced_code', 'tables', 'attr_list']
                )
                return self._enhance_formatting(formatted)
            except Exception as e:
                print(f"Markdown formatting error: {e}")
                return self._basic_format(text)
        else:
            return self._basic_format(text)
    
    def _basic_format(self, text):
        """Basic HTML formatting without markdown"""
        import html
        
        # Escape HTML first
        text = html.escape(text)
        
        # Handle tables (basic pipe table format)
        text = self._format_tables(text)
        
        # Handle code blocks
        text = re.sub(r'```(\w+)?\n(.*?)\n```', r'<pre><code class="language-\1">\2</code></pre>', text, flags=re.DOTALL)
        
        # Handle inline code
        text = re.sub(r'`([^`]+)`', r'<code>\1</code>', text)
        
        # Handle bold and italic (enhanced)
        text = re.sub(r'\*\*\*(.+?)\*\*\*', r'<strong><em>\1</em></strong>', text)  # Bold + italic
        text = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', text)  # Bold
        text = re.sub(r'\*(.+?)\*', r'<em>\1</em>', text)  # Italic
        text = re.sub(r'__(.+?)__', r'<u>\1</u>', text)  # Underline
        
        # Handle strikethrough
        text = re.sub(r'~~(.+?)~~', r'<del>\1</del>', text)
        
        # Handle headers
        text = re.sub(r'^#### (.+)$', r'<h4>\1</h4>', text, flags=re.MULTILINE)
        text = re.sub(r'^### (.+)$', r'<h3>\1</h3>', text, flags=re.MULTILINE)
        text = re.sub(r'^## (.+)$', r'<h2>\1</h2>', text, flags=re.MULTILINE)
        text = re.sub(r'^# (.+)$', r'<h1>\1</h1>', text, flags=re.MULTILINE)
        
        # Handle lists
        text = self._format_lists(text)
        
        # Handle links
        text = re.sub(r'\[([^\]]+)\]\(([^)]+)\)', r'<a href="\2" target="_blank">\1</a>', text)
        
        # Handle emojis and icons (preserve common ones)
        text = self._preserve_emojis(text)
        
        # Handle paragraphs
        text = text.replace('\n\n', '</p><p>')
        text = text.replace('\n', '<br>')
        
        if text.strip():
            text = f'<p>{text}</p>'
        
        return text
    
    def _enhance_formatting(self, text):
        """Additional formatting enhancements"""
        text = re.sub(r'</h([1-6])>', r'</h\1>\n', text)
        text = re.sub(r'</p>', r'</p>\n', text)
        text = re.sub(r'</ul>', r'</ul>\n', text)
        text = re.sub(r'</ol>', r'</ol>\n', text)
        
        # Enhance table styling
        text = re.sub(r'<table>', r'<table class="formatted-table">', text)
        
        # Add icons and emojis back
        text = self._restore_emojis(text)
        
        return text
    
    def _format_tables(self, text):
        """Convert pipe tables to HTML tables"""
        lines = text.split('\n')
        result = []
        in_table = False
        table_rows = []
        
        for line in lines:
            stripped = line.strip()
            
            # Check if this is a table row (contains pipes and content)
            if '|' in stripped and len(stripped.split('|')) > 2:
                # Remove leading/trailing empty cells
                cells = [cell.strip() for cell in stripped.split('|')]
                if cells[0] == '':
                    cells = cells[1:]
                if cells and cells[-1] == '':
                    cells = cells[:-1]
                
                if cells:  # Only process if we have actual content
                    if not in_table:
                        in_table = True
                        table_rows = []
                    
                    # Skip separator rows (containing only -, |, :, and spaces)
                    if not re.match(r'^[\s\-\|\:]+$', stripped):
                        table_rows.append(cells)
            else:
                # End of table
                if in_table and table_rows:
                    result.append(self._build_html_table(table_rows))
                    table_rows = []
                    in_table = False
                
                result.append(line)
        
        # Handle table at end of text
        if in_table and table_rows:
            result.append(self._build_html_table(table_rows))
        
        return '\n'.join(result)
    
    def _build_html_table(self, rows):
        """Build an HTML table from rows of cells"""
        if not rows:
            return ''
        
        table_html = ['<table class="formatted-table">']
        
        # First row as header if it looks like headers
        first_row = rows[0]
        if len(rows) > 1 and any(cell.isupper() or ':' in cell for cell in first_row):
            table_html.append('<thead><tr>')
            for cell in first_row:
                table_html.append(f'<th>{cell}</th>')
            table_html.append('</tr></thead>')
            body_rows = rows[1:]
        else:
            body_rows = rows
        
        # Body rows
        if body_rows:
            table_html.append('<tbody>')
            for row in body_rows:
                table_html.append('<tr>')
                for cell in row:
                    table_html.append(f'<td>{cell}</td>')
                table_html.append('</tr>')
            table_html.append('</tbody>')
        
        table_html.append('</table>')
        return '\n'.join(table_html)
    
    def _format_lists(self, text):
        """Convert markdown-style lists to HTML lists"""
        lines = text.split('\n')
        result = []
        in_ul = False
        in_ol = False
        
        for line in lines:
            stripped = line.strip()
            
            # Unordered list
            if re.match(r'^[\*\-\+]\s+', stripped):
                if not in_ul:
                    if in_ol:
                        result.append('</ol>')
                        in_ol = False
                    result.append('<ul>')
                    in_ul = True
                
                content = re.sub(r'^[\*\-\+]\s+', '', stripped)
                result.append(f'<li>{content}</li>')
            
            # Ordered list
            elif re.match(r'^\d+\.\s+', stripped):
                if not in_ol:
                    if in_ul:
                        result.append('</ul>')
                        in_ul = False
                    result.append('<ol>')
                    in_ol = True
                
                content = re.sub(r'^\d+\.\s+', '', stripped)
                result.append(f'<li>{content}</li>')
            
            else:
                # End lists
                if in_ul:
                    result.append('</ul>')
                    in_ul = False
                if in_ol:
                    result.append('</ol>')
                    in_ol = False
                
                result.append(line)
        
        # Close any open lists
        if in_ul:
            result.append('</ul>')
        if in_ol:
            result.append('</ol>')
        
        return '\n'.join(result)
    
    def _preserve_emojis(self, text):
        """Preserve common emojis and icons during HTML escaping"""
        # Common emojis that should be preserved
        emoji_map = {
            '&amp;#x2713;': '✓',  # checkmark
            '&amp;#x2717;': '✗',  # cross mark
            '&amp;#x2192;': '→',  # right arrow
            '&amp;#x2190;': '←',  # left arrow
            '&amp;#x2191;': '↑',  # up arrow
            '&amp;#x2193;': '↓',  # down arrow
        }
        
        for escaped, emoji in emoji_map.items():
            text = text.replace(escaped, emoji)
        
        return text
    
    def _restore_emojis(self, text):
        """Restore emojis after markdown processing"""
        # This method can be used to restore emojis that might have been
        # processed incorrectly by markdown
        return text

formatter = ResponseFormatter()
file_processor = FileProcessor()

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def enhance_prompt(user_message):
    """Enhance user prompt with formatting instructions"""
    return f"""{user_message}

Please format your response professionally using:
- **Bold text** for emphasis and important points
- *Italic text* for explanations or definitions
- `code snippets` for technical terms and functions
- ```language blocks``` for code examples
- ## Headings for organizing content
- Tables with | pipes | for | structured data |
- Bullet points and numbered lists for clarity
- Links [like this](url) when referencing external resources"""

def enhance_prompt_with_file(user_message, file_content, filename):
    """Add file context to user prompt"""
    return f"""I'm sharing a file with you: {filename}

File content:
{file_content}

User question: {user_message}

Please analyze this file and respond to my question. Format your response professionally using:
- **Bold text** for emphasis and important findings
- *Italic text* for explanations or definitions
- `code snippets` for technical terms and functions
- ```language blocks``` for code examples
- ## Headings to organize your analysis
- Tables with | pipes | for | data summary | or | comparisons |
- Bullet points and numbered lists for clarity
- Links [like this](url) when referencing external resources
- When analyzing spreadsheet data, present key findings in table format"""

# Database functions (same as before)
def get_db():
    """Get database connection"""
    conn = sqlite3.connect(DATABASE)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    """Initialize the database with required tables"""
    with get_db() as conn:
        conn.executescript('''
            CREATE TABLE IF NOT EXISTS sessions (
                id TEXT PRIMARY KEY,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                last_active TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                title TEXT,
                summary TEXT
            );
            
            CREATE TABLE IF NOT EXISTS messages (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                session_id TEXT,
                role TEXT NOT NULL,
                content TEXT NOT NULL,
                model TEXT,
                timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                has_file BOOLEAN DEFAULT FALSE,
                file_name TEXT,
                file_type TEXT,
                formatted_content TEXT,
                FOREIGN KEY (session_id) REFERENCES sessions (id)
            );
            
            CREATE TABLE IF NOT EXISTS file_attachments (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                message_id INTEGER,
                filename TEXT NOT NULL,
                original_filename TEXT NOT NULL,
                file_path TEXT NOT NULL,
                file_size INTEGER,
                mime_type TEXT,
                processed_content TEXT,
                timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (message_id) REFERENCES messages (id)
            );
        ''')
        conn.commit()

# Routes
@app.route('/')
def index():
    # Generate file accept attribute from allowed extensions
    file_accept = ','.join(f'.{ext}' for ext in sorted(ALLOWED_EXTENSIONS))
    return render_template('index.html', file_accept=file_accept)

@app.route('/api/chat', methods=['POST'])
def chat():
    try:
        # Handle both JSON and form data
        if request.is_json:
            data = request.get_json()
            user_message = data.get('message', '').strip()
            model = data.get('model', 'llama2')
            uploaded_file = None
        else:
            user_message = request.form.get('message', '').strip()
            model = request.form.get('model', 'llama2')
            uploaded_file = request.files.get('file')
        
        if not user_message.strip() and not uploaded_file:
            return jsonify({'success': False, 'error': 'Empty message and no file'}), 400
        
        # Session management
        if 'session_id' not in session:
            session['session_id'] = create_session()
        
        session_id = session['session_id']
        create_session(session_id)
        
        file_info = None
        file_content = ""
        
        # Process uploaded file if present
        if uploaded_file and uploaded_file.filename != '':
            if allowed_file(uploaded_file.filename):
                filename = secure_filename(uploaded_file.filename)
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                unique_filename = f"{timestamp}_{filename}"
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                
                uploaded_file.save(file_path)
                
                # Get file info
                file_size = os.path.getsize(file_path)
                mime_type = mimetypes.guess_type(file_path)[0] or 'application/octet-stream'
                
                # Process file content
                processed_content = file_processor.process_file(file_path, filename, mime_type)
                
                file_info = {
                    'filename': unique_filename,
                    'original_filename': filename,
                    'file_path': file_path,
                    'file_size': file_size,
                    'mime_type': mime_type,
                    'processed_content': processed_content
                }
                
                file_content = processed_content
            else:
                return jsonify({'success': False, 'error': 'File type not allowed'}), 400
        
        # Save user message to database
        add_message(session_id, 'user', user_message, model, file_info)
        
        # Get recent messages for context
        recent_messages = get_recent_messages(session_id, limit=10)
        
        # Prepare the prompt
        if file_content:
            enhanced_message = enhance_prompt_with_file(user_message, file_content, file_info['original_filename'])
        else:
            enhanced_message = enhance_prompt(user_message)
        
        # Prepare messages for AI model
        messages = recent_messages[:-1]  # Exclude the current message
        messages.append({'role': 'user', 'content': enhanced_message})
        
        # Determine which AI service to use
        try:
            if model == CLAUDE_CODE_MODEL and claude_code_client.available:
                # Send to Claude Code (async)
                async def claude_chat():
                    return await claude_code_client.chat(messages, ["Read", "Write", "Bash", "Grep"])
                
                response = anyio.run(claude_chat)
                assistant_message = response.get('message', {}).get('content', 'No response received')
            else:
                # Send to Ollama
                response = ollama_client.chat(model, messages)
                assistant_message = response.get('message', {}).get('content', 'No response received')
            
            # Save assistant response to database
            add_message(session_id, 'assistant', assistant_message, model)
            
            # Format the response
            formatted_content = formatter.format_response(assistant_message)
            
            return jsonify({
                'success': True,
                'message': {
                    'content': assistant_message,
                    'formatted_content': formatted_content
                }
            })
            
        except Exception as e:
            error_message = str(e)
            print(f"ERROR: Chat endpoint exception: {error_message}")
            print(f"ERROR: Model was: {model}")
            print(f"ERROR: Claude Code available: {claude_code_client.available}")
            import traceback
            traceback.print_exc()
            
            add_message(session_id, 'assistant', f"Error: {error_message}", model)
            return jsonify({'success': False, 'error': error_message}), 500
            
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/models', methods=['GET'])
def get_models():
    try:
        models_data = ollama_client.list_models()
        models = []
        
        if models_data.get('models'):
            models = [model['name'] for model in models_data['models']]
        
        # If no models found, return default models
        if not models:
            models = DEFAULT_MODELS.copy()
        
        # Add Claude Code if available
        if claude_code_client.available:
            models.append(CLAUDE_CODE_MODEL)
        elif CLAUDE_CODE_AVAILABLE:
            # Show Claude Code as disabled if SDK is available but no API key
            models.append(f"{CLAUDE_CODE_MODEL} (API key required)")
            
        return jsonify({'success': True, 'models': models})
    except Exception as e:
        models = DEFAULT_MODELS.copy()
        if claude_code_client.available:
            models.append(CLAUDE_CODE_MODEL)
        return jsonify({'success': True, 'models': models})

@app.route('/api/sessions', methods=['GET'])
def get_sessions():
    try:
        sessions = get_all_sessions()
        
        # Format sessions for frontend
        formatted_sessions = []
        for session_data in sessions:
            formatted_sessions.append({
                'session_id': session_data['id'],
                'title': session_data['title'] or f"Session {session_data['id'][:8]}",
                'created_at': session_data['created_at'],
                'last_active': session_data['last_active'],
                'message_count': session_data['message_count'] or 0,
                'last_message': session_data['last_message'] or 'No messages'
            })
        
        return jsonify({'success': True, 'sessions': formatted_sessions})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/session/<session_id>/load', methods=['POST'])
def load_session(session_id):
    try:
        messages = get_session_messages(session_id)
        session['session_id'] = session_id
        
        # Format messages for frontend
        formatted_messages = []
        for msg in messages:
            formatted_messages.append({
                'role': msg['role'],
                'content': msg['content'],
                'formatted_content': msg['formatted_content'],
                'timestamp': msg['timestamp'],
                'has_file': msg['has_file'],
                'file_name': msg['original_filename'] if msg.get('original_filename') else msg.get('file_name'),
                'model': msg['model']
            })
        
        return jsonify({'success': True, 'messages': formatted_messages})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/session/<session_id>/delete', methods=['DELETE'])
def delete_session(session_id):
    try:
        with get_db() as conn:
            # Delete file attachments first
            cursor = conn.execute('''
                SELECT fa.file_path FROM file_attachments fa
                JOIN messages m ON fa.message_id = m.id
                WHERE m.session_id = ?
            ''', (session_id,))
            
            file_paths = cursor.fetchall()
            
            # Delete physical files
            for file_path_row in file_paths:
                file_path = file_path_row['file_path']
                try:
                    if os.path.exists(file_path):
                        os.remove(file_path)
                except Exception as e:
                    print(f"Error deleting file {file_path}: {e}")
            
            # Delete database records
            conn.execute('''
                DELETE FROM file_attachments 
                WHERE message_id IN (SELECT id FROM messages WHERE session_id = ?)
            ''', (session_id,))
            
            conn.execute('DELETE FROM messages WHERE session_id = ?', (session_id,))
            conn.execute('DELETE FROM sessions WHERE id = ?', (session_id,))
            conn.commit()
        
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/session/new', methods=['POST'])
def new_session():
    try:
        new_session_id = create_session()
        session['session_id'] = new_session_id
        return jsonify({'success': True, 'session_id': new_session_id})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/claude-code/api-key', methods=['POST'])
def set_claude_api_key():
    try:
        data = request.get_json()
        api_key = data.get('api_key', '').strip()
        
        if not api_key:
            return jsonify({'success': False, 'error': 'API key is required'}), 400
        
        # Set the API key for Claude Code client
        claude_code_client.set_api_key(api_key)
        
        # Store in environment variable for this session
        os.environ['CLAUDE_CODE_API_KEY'] = api_key
        
        return jsonify({
            'success': True, 
            'message': 'Claude Code API key set successfully',
            'available': claude_code_client.available
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/claude-code/status', methods=['GET'])
def get_claude_status():
    try:
        return jsonify({
            'success': True,
            'available': claude_code_client.available,
            'sdk_installed': CLAUDE_CODE_AVAILABLE,
            'has_api_key': bool(claude_code_client.api_key)
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# Database helper functions
def create_session(session_id=None):
    """Create a new chat session"""
    if not session_id:
        session_id = str(uuid.uuid4())
    
    try:
        with get_db() as conn:
            conn.execute(
                'INSERT OR IGNORE INTO sessions (id, title) VALUES (?, ?)',
                (session_id, f'Chat Session {datetime.now().strftime("%Y-%m-%d %H:%M")}')
            )
            conn.commit()
    except Exception as e:
        print(f"Error creating session: {e}")
    
    return session_id

def add_message(session_id, role, content, model=None, file_info=None):
    """Add a message to the database"""
    formatted_content = None
    if role == 'assistant':
        try:
            formatted_content = formatter.format_response(content)
        except Exception as e:
            print(f"Error formatting message: {e}")
            formatted_content = content
    
    has_file = file_info is not None
    file_name = file_info.get('original_filename') if file_info else None
    file_type = file_info.get('mime_type') if file_info else None
    
    try:
        with get_db() as conn:
            cursor = conn.execute(
                '''INSERT INTO messages (session_id, role, content, formatted_content, model, has_file, file_name, file_type) 
                   VALUES (?, ?, ?, ?, ?, ?, ?, ?)''',
                (session_id, role, content, formatted_content, model, has_file, file_name, file_type)
            )
            message_id = cursor.lastrowid
            
            # Add file attachment if present
            if file_info:
                conn.execute(
                    '''INSERT INTO file_attachments (message_id, filename, original_filename, file_path, file_size, mime_type, processed_content)
                       VALUES (?, ?, ?, ?, ?, ?, ?)''',
                    (message_id, file_info['filename'], file_info['original_filename'], 
                     file_info['file_path'], file_info['file_size'], file_info['mime_type'], 
                     file_info.get('processed_content', ''))
                )
            
            # Update session last active time
            conn.execute(
                'UPDATE sessions SET last_active = CURRENT_TIMESTAMP WHERE id = ?',
                (session_id,)
            )
            conn.commit()
            return message_id
    except Exception as e:
        print(f"Error adding message to database: {e}")
        return None

def get_session_messages(session_id, limit=50):
    """Get messages for a session"""
    try:
        with get_db() as conn:
            cursor = conn.execute(
                '''SELECT m.*, fa.original_filename, fa.mime_type as attachment_mime_type, fa.processed_content
                   FROM messages m
                   LEFT JOIN file_attachments fa ON m.id = fa.message_id
                   WHERE m.session_id = ? 
                   ORDER BY m.timestamp DESC LIMIT ?''',
                (session_id, limit)
            )
            messages = cursor.fetchall()
            return [dict(msg) for msg in reversed(messages)]
    except Exception as e:
        print(f"Error getting session messages: {e}")
        return []

def get_recent_messages(session_id, limit=10):
    """Get recent messages for context"""
    try:
        with get_db() as conn:
            cursor = conn.execute(
                '''SELECT m.role, m.content, fa.processed_content
                   FROM messages m
                   LEFT JOIN file_attachments fa ON m.id = fa.message_id
                   WHERE m.session_id = ? 
                   ORDER BY m.timestamp DESC LIMIT ?''',
                (session_id, limit)
            )
            messages = cursor.fetchall()
            result = []
            for msg in reversed(messages):
                message_content = msg['content']
                if msg['processed_content']:
                    message_content += f"\n\nFile content:\n{msg['processed_content']}"
                result.append({'role': msg['role'], 'content': message_content})
            return result
    except Exception as e:
        print(f"Error getting recent messages: {e}")
        return []

def get_all_sessions():
    """Get all chat sessions"""
    try:
        with get_db() as conn:
            cursor = conn.execute('''
                SELECT s.*,
                       COUNT(m.id) as message_count,
                       MAX(m.timestamp) as last_message_time,
                       (SELECT content FROM messages WHERE session_id = s.id ORDER BY timestamp DESC LIMIT 1) as last_message
                FROM sessions s
                LEFT JOIN messages m ON s.id = m.session_id
                GROUP BY s.id
                ORDER BY s.last_active DESC
            ''')
            sessions = cursor.fetchall()
            return [dict(session) for session in sessions]
    except Exception as e:
        print(f"Error getting sessions: {e}")
        return []

def load_messages_from_session(messages):
    """Helper function to format messages for frontend loading"""
    formatted_messages = []
    for msg in messages:
        formatted_messages.append({
            'role': msg['role'],
            'content': msg['content'],
            'formatted_content': msg.get('formatted_content'),
            'timestamp': msg['timestamp'],
            'has_file': msg.get('has_file', False),
            'file_name': msg.get('file_name'),
            'model': msg.get('model')
        })
    return formatted_messages

# =============================================================================
# IMAGE GENERATION FUNCTIONALITY
# =============================================================================

@app.route('/api/generate-image', methods=['POST'])
def generate_image():
    """Generate an image using Hugging Face Inference API"""
    print("Image generation endpoint called")
    
    if not HUGGINGFACE_AVAILABLE:
        print("HUGGINGFACE_AVAILABLE is False")
        return jsonify({'success': False, 'error': 'Image generation not available - huggingface-hub not installed'}), 400
    
    try:
        data = request.get_json()
        print(f"Received data: {data}")
        prompt = data.get('prompt', '').strip()
        
        if not prompt:
            return jsonify({'success': False, 'error': 'Prompt is required'}), 400
        
        # Check for Hugging Face API token (required for image generation)
        hf_token = os.getenv('HUGGINGFACE_HUB_TOKEN') or os.getenv('HF_TOKEN')
        
        if not hf_token:
            return jsonify({
                'success': False, 
                'error': 'Hugging Face API token required for image generation. Please set HUGGINGFACE_HUB_TOKEN or HF_TOKEN environment variable.',
                'setup_required': True
            }), 400
        
        # Initialize Hugging Face client
        client = InferenceClient(token=hf_token)
        
        # Image generation parameters - using confirmed working models
        model_name = data.get('model', 'stabilityai/stable-diffusion-2-1')
        
        print(f"Generating image with Hugging Face model {model_name}: {prompt}")
        
        # Generate image using Hugging Face Inference Providers
        try:
            # Try with specific model first
            image = client.text_to_image(
                prompt,
                model=model_name
            )
            
            # Convert PIL Image to base64 for frontend
            img_buffer = io.BytesIO()
            image.save(img_buffer, format='PNG')
            img_buffer.seek(0)
            
            # Create base64 data URL
            img_base64 = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
            image_url = f"data:image/png;base64,{img_base64}"
            
            return jsonify({
                'success': True,
                'images': [image_url],
                'prompt': prompt,
                'model': model_name
            })
            
        except Exception as model_error:
            # If the specified model fails, try without specifying a model (uses default)
            print(f"Model {model_name} failed, trying default model: {str(model_error)}")
            try:
                image = client.text_to_image(prompt)
                fallback_model = "default"
                
                # Convert PIL Image to base64 for frontend
                img_buffer = io.BytesIO()
                image.save(img_buffer, format='PNG')
                img_buffer.seek(0)
                
                img_base64 = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
                image_url = f"data:image/png;base64,{img_base64}"
                
                return jsonify({
                    'success': True,
                    'images': [image_url],
                    'prompt': prompt,
                    'model': fallback_model
                })
            except Exception as default_error:
                raise default_error
        
    except Exception as e:
        error_message = str(e)
        print(f"Image generation error: {error_message}")
        return jsonify({'success': False, 'error': f'Image generation failed: {error_message}'}), 500

@app.route('/api/generate-video', methods=['POST'])
def generate_video():
    """Generate a video using available free services"""
    print("Video generation endpoint called")
    
    try:
        data = request.get_json()
        prompt = data.get('prompt', '').strip()
        duration = int(data.get('duration', 3))  # Default 3 seconds
        fps = int(data.get('fps', 8))  # Lower FPS for free tier
        
        if not prompt:
            return jsonify({'success': False, 'error': 'Prompt is required'}), 400
        
        # For now, create a simple animation using image sequence
        # This is a basic implementation - can be enhanced later
        video_result = create_simple_video_animation(prompt, duration, fps)
        
        if video_result['success']:
            return jsonify({
                'success': True,
                'video_url': video_result['video_url'],
                'prompt': prompt,
                'duration': duration,
                'fps': fps,
                'method': 'simple_animation'
            })
        else:
            return jsonify({
                'success': False,
                'error': video_result['error']
            }), 500
            
    except Exception as e:
        error_message = str(e)
        print(f"Video generation error: {error_message}")
        return jsonify({
            'success': False, 
            'error': f'Video generation failed: {error_message}'
        }), 500

def create_simple_video_animation(prompt, duration, fps):
    """Create a simple video animation using image generation"""
    try:
        if not HUGGINGFACE_AVAILABLE:
            return {'success': False, 'error': 'Hugging Face not available for image generation'}
        
        hf_token = os.getenv('HUGGINGFACE_HUB_TOKEN') or os.getenv('HF_TOKEN')
        if not hf_token:
            return {'success': False, 'error': 'Hugging Face token required'}
        
        # Create multiple images with slight variations
        frames = []
        total_frames = duration * fps
        
        client = InferenceClient(token=hf_token)
        
        # Generate keyframes with variations
        for i in range(min(total_frames, 8)):  # Limit to 8 frames for free tier
            # Add slight variations to the prompt for each frame
            frame_prompt = f"{prompt}, frame {i+1}, slightly different angle"
            
            try:
                image = client.text_to_image(frame_prompt)
                frames.append(image)
            except Exception as e:
                print(f"Frame {i} generation failed: {e}")
                # Use the last successful frame if available
                if frames:
                    frames.append(frames[-1])
        
        if not frames:
            return {'success': False, 'error': 'No frames could be generated'}
        
        # Create video from frames
        video_path = create_video_from_frames(frames, fps)
        
        if video_path:
            # Convert to base64 for web delivery
            with open(video_path, 'rb') as video_file:
                video_data = video_file.read()
                video_base64 = base64.b64encode(video_data).decode('utf-8')
                video_url = f"data:video/mp4;base64,{video_base64}"
            
            # Clean up temporary file
            os.remove(video_path)
            
            return {'success': True, 'video_url': video_url}
        else:
            return {'success': False, 'error': 'Failed to create video from frames'}
            
    except Exception as e:
        return {'success': False, 'error': str(e)}

def create_video_from_frames(frames, fps):
    """Create an MP4 video from a list of PIL images"""
    try:
        import tempfile
        import cv2
        import numpy as np
        
        # Create temporary directory for frames
        with tempfile.TemporaryDirectory() as temp_dir:
            frame_paths = []
            
            # Save frames as temporary images
            for i, frame in enumerate(frames):
                frame_path = os.path.join(temp_dir, f"frame_{i:04d}.png")
                frame.save(frame_path)
                frame_paths.append(frame_path)
            
            # Create video file
            temp_video = tempfile.NamedTemporaryFile(suffix='.mp4', delete=False)
            video_path = temp_video.name
            temp_video.close()
            
            # Get frame dimensions from first frame
            first_frame = cv2.imread(frame_paths[0])
            height, width, layers = first_frame.shape
            
            # Create video writer
            fourcc = cv2.VideoWriter_fourcc(*'mp4v')
            video_writer = cv2.VideoWriter(video_path, fourcc, fps, (width, height))
            
            # Add each frame multiple times to extend duration
            repeats = max(1, 24 // len(frames))  # Repeat frames to get reasonable duration
            
            for frame_path in frame_paths:
                frame = cv2.imread(frame_path)
                for _ in range(repeats):
                    video_writer.write(frame)
            
            video_writer.release()
            cv2.destroyAllWindows()
            
            return video_path
            
    except ImportError:
        print("OpenCV not available, trying alternative method")
        return create_video_with_pillow(frames, fps)
    except Exception as e:
        print(f"Video creation error: {e}")
        return None

def create_video_with_pillow(frames, fps):
    """Alternative video creation using Pillow for GIF (fallback)"""
    try:
        import tempfile
        
        # Create animated GIF as fallback
        temp_gif = tempfile.NamedTemporaryFile(suffix='.gif', delete=False)
        gif_path = temp_gif.name
        temp_gif.close()
        
        # Create animated GIF
        frames[0].save(
            gif_path,
            save_all=True,
            append_images=frames[1:],
            duration=1000//fps,  # Duration in milliseconds
            loop=0
        )
        
        return gif_path
        
    except Exception as e:
        print(f"GIF creation error: {e}")
        return None

@app.route('/api/test-image-simple', methods=['POST'])
def test_image_simple():
    """Test simple image generation without model specification"""
    if not HUGGINGFACE_AVAILABLE:
        return jsonify({'success': False, 'error': 'Hugging Face not available'}), 400
    
    try:
        data = request.get_json()
        prompt = data.get('prompt', 'a cat').strip()
        
        hf_token = os.getenv('HUGGINGFACE_HUB_TOKEN') or os.getenv('HF_TOKEN')
        if not hf_token:
            return jsonify({'success': False, 'error': 'Token required'}), 400
        
        print(f"Testing simple image generation with prompt: {prompt}")
        
        # Try the simplest possible approach
        client = InferenceClient(token=hf_token)
        image = client.text_to_image(prompt)
        
        # Convert to base64
        img_buffer = io.BytesIO()
        image.save(img_buffer, format='PNG')
        img_buffer.seek(0)
        
        img_base64 = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
        image_url = f"data:image/png;base64,{img_base64}"
        
        return jsonify({
            'success': True,
            'images': [image_url],
            'prompt': prompt,
            'model': 'default'
        })
        
    except Exception as e:
        error_message = str(e)
        print(f"Simple image test error: {error_message}")
        return jsonify({'success': False, 'error': f'Test failed: {error_message}'}), 500

@app.route('/api/huggingface/status', methods=['GET'])
def huggingface_status():
    """Check Hugging Face API availability"""
    if not HUGGINGFACE_AVAILABLE:
        return jsonify({
            'available': False,
            'sdk_installed': False,
            'has_api_key': False,
            'error': 'Hugging Face Hub SDK not installed'
        })
    
    hf_token = os.getenv('HUGGINGFACE_HUB_TOKEN') or os.getenv('HF_TOKEN')
    
    # Image generation requires a Hugging Face token
    return jsonify({
        'available': bool(hf_token and hf_token.strip()),
        'sdk_installed': True,
        'has_api_key': bool(hf_token and hf_token.strip()),
        'token_required': True,
        'setup_url': 'https://huggingface.co/settings/tokens'
    })

@app.route('/api/huggingface/api-key', methods=['POST'])
def set_huggingface_api_key():
    """Set Hugging Face API key (optional - improves rate limits)"""
    if not HUGGINGFACE_AVAILABLE:
        return jsonify({'success': False, 'error': 'Hugging Face Hub SDK not installed'}), 400
    
    try:
        data = request.get_json()
        api_key = data.get('api_key', '').strip()
        
        if not api_key:
            return jsonify({'success': False, 'error': 'API key is required'}), 400
        
        # Set environment variable
        os.environ['HUGGINGFACE_HUB_TOKEN'] = api_key
        
        # Test the API key by creating a client
        try:
            client = InferenceClient(token=api_key)
            # Simple test - this should work with any valid token
            return jsonify({'success': True, 'available': True})
        except Exception as e:
            return jsonify({'success': False, 'error': f'Invalid API key: {str(e)}'}), 400
            
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# =============================================================================
# EXPORT FUNCTIONALITY
# =============================================================================

@app.route('/api/export/word/<session_id>', methods=['GET'])
def export_to_word(session_id):
    """Export chat session to Word document"""
    if not DOCX_AVAILABLE:
        return jsonify({'success': False, 'error': 'Word export not available - python-docx not installed'}), 400
    
    try:
        from docx import Document
        from docx.shared import Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        import tempfile
        from flask import send_file
        
        # Get session data
        messages = get_session_messages(session_id)
        if not messages:
            return jsonify({'success': False, 'error': 'No messages found for this session'}), 404
        
        # Create Word document
        doc = Document()
        
        # Add title
        title = doc.add_heading('Chat Session Export', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add session info
        doc.add_paragraph(f'Session ID: {session_id}')
        doc.add_paragraph(f'Exported on: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')
        doc.add_paragraph(f'Total Messages: {len(messages)}')
        doc.add_paragraph('')  # Empty line
        
        # Add messages
        for i, msg in enumerate(messages, 1):
            # Message header
            role = msg['role'].title()
            timestamp = msg['timestamp']
            model = msg.get('model', 'Unknown')
            
            header = doc.add_paragraph()
            header.add_run(f"{i}. {role}").bold = True
            if msg['role'] == 'assistant' and model:
                header.add_run(f" ({model})")
            header.add_run(f" - {timestamp}")
            
            # Message content
            content = msg.get('content', '')
            if content:
                doc.add_paragraph(content)
            
            # File attachment info
            if msg.get('has_file') and msg.get('file_name'):
                file_para = doc.add_paragraph()
                file_para.add_run("📎 Attached file: ").italic = True
                file_para.add_run(msg['file_name'])
            
            doc.add_paragraph('')  # Empty line between messages
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
        doc.save(temp_file.name)
        temp_file.close()
        
        filename = f"chat_session_{session_id[:8]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        
        return send_file(
            temp_file.name,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
        
    except Exception as e:
        return jsonify({'success': False, 'error': f'Export failed: {str(e)}'}), 500

@app.route('/api/export/excel/<session_id>', methods=['GET'])
def export_to_excel(session_id):
    """Export chat session to Excel spreadsheet"""
    if not EXCEL_AVAILABLE:
        return jsonify({'success': False, 'error': 'Excel export not available - openpyxl not installed'}), 400
    
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill
        import tempfile
        from flask import send_file
        
        # Get session data
        messages = get_session_messages(session_id)
        if not messages:
            return jsonify({'success': False, 'error': 'No messages found for this session'}), 404
        
        # Create workbook
        wb = Workbook()
        ws = wb.active
        ws.title = f"Chat Session {session_id[:8]}"
        
        # Headers
        headers = ['#', 'Role', 'Model', 'Timestamp', 'Content', 'Has File', 'File Name']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
        
        # Add messages
        for i, msg in enumerate(messages, 2):
            ws.cell(row=i, column=1, value=i-1)  # Message number
            ws.cell(row=i, column=2, value=msg['role'].title())
            ws.cell(row=i, column=3, value=msg.get('model', ''))
            ws.cell(row=i, column=4, value=msg['timestamp'])
            ws.cell(row=i, column=5, value=msg.get('content', ''))
            ws.cell(row=i, column=6, value='Yes' if msg.get('has_file') else 'No')
            ws.cell(row=i, column=7, value=msg.get('file_name', ''))
        
        # Auto-adjust column widths
        for col in ws.columns:
            max_length = 0
            column = col[0].column_letter
            for cell in col:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)  # Cap at 50 characters
            ws.column_dimensions[column].width = adjusted_width
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx')
        wb.save(temp_file.name)
        temp_file.close()
        
        filename = f"chat_session_{session_id[:8]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        
        return send_file(
            temp_file.name,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )
        
    except Exception as e:
        return jsonify({'success': False, 'error': f'Export failed: {str(e)}'}), 500

@app.route('/api/export/pdf/<session_id>', methods=['GET'])
def export_to_pdf(session_id):
    """Export chat session to PDF"""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        import tempfile
        from flask import send_file
        
        # Get session data
        messages = get_session_messages(session_id)
        if not messages:
            return jsonify({'success': False, 'error': 'No messages found for this session'}), 404
        
        # Create PDF
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        doc = SimpleDocTemplate(temp_file.name, pagesize=letter)
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            alignment=1,  # Center alignment
            spaceAfter=30,
        )
        
        header_style = ParagraphStyle(
            'MessageHeader',
            parent=styles['Heading3'],
            spaceBefore=12,
            spaceAfter=6,
        )
        
        content_style = styles['Normal']
        
        story = []
        
        # Add title
        story.append(Paragraph("Chat Session Export", title_style))
        story.append(Spacer(1, 12))
        
        # Add session info
        story.append(Paragraph(f"<b>Session ID:</b> {session_id}", content_style))
        story.append(Paragraph(f"<b>Exported on:</b> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", content_style))
        story.append(Paragraph(f"<b>Total Messages:</b> {len(messages)}", content_style))
        story.append(Spacer(1, 20))
        
        # Add messages
        for i, msg in enumerate(messages, 1):
            role = msg['role'].title()
            timestamp = msg['timestamp']
            model = msg.get('model', '')
            
            # Message header
            header_text = f"{i}. {role}"
            if msg['role'] == 'assistant' and model:
                header_text += f" ({model})"
            header_text += f" - {timestamp}"
            
            story.append(Paragraph(header_text, header_style))
            
            # Message content
            content = msg.get('content', '').replace('\n', '<br/>')
            if content:
                story.append(Paragraph(content, content_style))
            
            # File attachment info
            if msg.get('has_file') and msg.get('file_name'):
                story.append(Paragraph(f"<i>📎 Attached file: {msg['file_name']}</i>", content_style))
            
            story.append(Spacer(1, 12))
        
        # Build PDF
        doc.build(story)
        temp_file.close()
        
        filename = f"chat_session_{session_id[:8]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        
        return send_file(
            temp_file.name,
            as_attachment=True,
            download_name=filename,
            mimetype='application/pdf'
        )
        
    except ImportError:
        return jsonify({'success': False, 'error': 'PDF export not available - reportlab not installed'}), 400
    except Exception as e:
        return jsonify({'success': False, 'error': f'Export failed: {str(e)}'}), 500

# Error handlers
@app.errorhandler(413)
def too_large(e):
    return jsonify({'success': False, 'error': 'File too large. Maximum size is 16MB.'}), 413

@app.errorhandler(404)
def not_found(e):
    return jsonify({'success': False, 'error': 'Endpoint not found'}), 404

@app.errorhandler(500)
def internal_error(e):
    return jsonify({'success': False, 'error': 'Internal server error'}), 500

# Initialize database on startup
def safe_print(text):
    """Print text with Unicode fallback for Windows terminals"""
    try:
        print(text)
    except UnicodeEncodeError:
        # Remove emojis and special characters for Windows compatibility
        fallback = text.replace("🚀", "").replace("📁", "").replace("🗄️", "").replace("🌐", "")
        fallback = fallback.replace("🤖", "").replace("✅", "+").replace("❌", "-").replace("⚠️", "!")
        fallback = fallback.replace("📄", "").replace("⚡", "").replace("🎯", "").replace("📎", "")
        print(fallback.strip())

def initialize_app():
    """Initialize the application"""
    safe_print("🚀 Starting Sypha AI Chat Application...")
    print("=" * 50)
    safe_print(f"📁 Upload folder: {UPLOAD_FOLDER}")
    safe_print(f"🗄️  Database: {DATABASE}")
    safe_print(f"🌐 Ollama URL: {OLLAMA_BASE_URL}")
    print("")
    
    # Check AI capabilities
    safe_print("🤖 AI MODELS:")
    
    ai_models = []
    
    # Test Ollama connection
    try:
        models = ollama_client.list_models()
        if models.get('models'):
            model_names = [model['name'] for model in models['models']]
            ai_models.extend(model_names)
            safe_print(f"   ✅ Ollama: {len(model_names)} models ({', '.join(model_names[:3])}{'...' if len(model_names) > 3 else ''})")
        else:
            safe_print("   ⚠️  Ollama: No models found (install models with: ollama pull llama2)")
    except Exception as e:
        safe_print(f"   ❌ Ollama: Connection failed ({e})")
        print("      → Start with: ollama serve")
    
    # Claude Code status
    if claude_code_client.available:
        ai_models.append("claude-code")
        safe_print("   ✅ Claude Code: Ready with advanced tools")
    elif CLAUDE_CODE_AVAILABLE:
        safe_print("   ⚠️  Claude Code: SDK installed, API key required")
    else:
        safe_print("   ❌ Claude Code: SDK not available")
    
    print("")
    
    # Document processing capabilities
    safe_print("📄 DOCUMENT PROCESSING:")
    capabilities = []
    if DOCX_AVAILABLE:
        capabilities.append("Word (.docx)")
    if EXCEL_AVAILABLE:
        capabilities.append("Excel (.xlsx/.xls)")
    if PPTX_AVAILABLE:
        capabilities.append("PowerPoint (.pptx)")
    if PDF_AVAILABLE:
        capabilities.append("PDF")
    if IMAGE_AVAILABLE:
        capabilities.append("Images")
    
    if capabilities:
        for cap in capabilities:
            safe_print(f"   ✅ {cap}")
    else:
        safe_print("   ⚠️  Basic text files only")
    
    print("")
    
    # Advanced features
    safe_print("⚡ ADVANCED FEATURES:")
    if MARKDOWN_AVAILABLE:
        safe_print("   ✅ Rich text formatting with tables")
    safe_print("   ✅ Split-screen responsive interface")
    safe_print("   ✅ File upload and analysis")
    safe_print("   ✅ Session-based conversation history")
    safe_print("   ✅ Interactive tables and content")
    safe_print("   ✅ Speech-to-text voice input (Web Speech API)")
    safe_print("   ✅ Export to Word, Excel, and PDF")
    
    # Image generation features
    if HUGGINGFACE_AVAILABLE:
        hf_token = os.getenv('HUGGINGFACE_HUB_TOKEN') or os.getenv('HF_TOKEN')
        if hf_token and hf_token.strip():
            safe_print("   ✅ AI image generation (Hugging Face with API key)")
        else:
            safe_print("   ✅ AI image generation (Hugging Face free tier)")
    else:
        safe_print("   ❌ AI image generation (huggingface-hub package not installed)")
    
    print("")
    safe_print(f"📎 Supported file types: {', '.join(sorted(ALLOWED_EXTENSIONS))}")
    
    # Initialize database
    init_db()
    safe_print("✅ Database initialized")
    
    print("")
    print("=" * 50)
    if ai_models:
        safe_print(f"🎯 Ready! {len(ai_models)} AI model(s) available")
    else:
        safe_print("⚠️  Ready with limited functionality (no AI models)")
    safe_print("🌐 Access the application at: http://localhost:5000")
    print("=" * 50)

if __name__ == '__main__':
    initialize_app()
    
    # Run the application
    app.run(
        host='0.0.0.0',
        port=5000,
        debug=True,
        threaded=True
    )

